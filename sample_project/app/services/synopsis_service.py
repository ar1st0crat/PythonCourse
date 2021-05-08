import os
import uuid
from datetime import datetime
import pdfplumber
import pptx
from app.models.synopsis import Synopsis
from app.exts import db


PROJECT_FOLDER = os.path.dirname(os.path.realpath(__file__))
UPLOAD_FOLDER = fr'{PROJECT_FOLDER}/../uploads'

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)


def extract_pdf(pdf_file, topic_id):

    url = f'{uuid.uuid1()}.pdf'
    filename = f'{UPLOAD_FOLDER}\\{url}'
    pdf_file.save(filename)

    full_text = []

    with pdfplumber.open(filename) as pdf:
        for page in pdf.pages:
            text = page.dedupe_chars().extract_text()
            full_text.append(text)
        
    full_text = '\n === \n'.join(full_text)

    synopsis = Synopsis()
    synopsis.created_at = datetime.utcnow()
    synopsis.topic_id = topic_id
    synopsis.text = full_text
    synopsis.file_url = url

    db.session.add(synopsis)
    db.session.commit()

    return full_text


def extract_pptx(pptx_file, topic_id):

    url = f'{uuid.uuid1()}.pptx'
    filename = f'{UPLOAD_FOLDER}\\{url}'
    pptx_file.save(filename)
    
    prs = pptx.Presentation(filename)

    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        text_runs.append('\n === \n')

    text = ' '.join(text_runs)

    while '  ' in text:
        text = text.replace('  ', ' ')

    synopsis = Synopsis()
    synopsis.created_at = datetime.utcnow()
    synopsis.topic_id = topic_id
    synopsis.text = text
    synopsis.file_url = url

    db.session.add(synopsis)
    db.session.commit()

    return text
